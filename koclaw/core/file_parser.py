import base64
import mimetypes
import subprocess
from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader

try:
    from docx import Document
except ImportError:
    Document = None  # type: ignore[assignment,misc]

from koclaw.core.prompt_guard import wrap_external_content

TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".xml", ".html", ".log"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
PDF_EXTENSIONS = {".pdf"}
HWPX_EXTENSIONS = {".hwpx"}
HWP_EXTENSIONS = {".hwp"}
DOCX_EXTENSIONS = {".docx"}
XLSX_EXTENSIONS = {".xlsx", ".xls"}
PPTX_EXTENSIONS = {".pptx"}
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def _rows_to_markdown(rows: list[list[str]]) -> str:
    """2D 셀 목록을 마크다운 테이블 문자열로 변환합니다."""
    if not rows:
        return ""
    header = rows[0]
    sep = ["---"] * len(header)
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(sep) + " |",
    ]
    for row in rows[1:]:
        # 열 수가 헤더보다 짧으면 빈 칸으로 채움
        padded = row + [""] * max(0, len(header) - len(row))
        lines.append("| " + " | ".join(padded[: len(header)]) + " |")
    return "\n".join(lines)


@dataclass
class ParsedFile:
    name: str
    content: str
    mime_type: str

    @property
    def is_image(self) -> bool:
        return self.mime_type.startswith("image/")

    def to_image_part(self) -> dict:
        return {"type": "image", "data": self.content, "mime_type": self.mime_type}

    def to_llm_context(self) -> str:
        return wrap_external_content(f"첨부파일: {self.name}", self.content)


class FileParser:
    async def parse(self, path: str | Path) -> ParsedFile:
        path = Path(path)
        ext = path.suffix.lower()

        if ext in TEXT_EXTENSIONS:
            return self._parse_text(path)
        if ext in PDF_EXTENSIONS:
            return self._parse_pdf(path)
        if ext in HWPX_EXTENSIONS:
            return self._parse_hwpx(path)
        if ext in HWP_EXTENSIONS:
            return self._parse_hwp(path)
        if ext in DOCX_EXTENSIONS:
            return self._parse_docx(path)
        if ext in XLSX_EXTENSIONS:
            return self._parse_xlsx(path)
        if ext in PPTX_EXTENSIONS:
            return self._parse_pptx(path)
        if ext in IMAGE_EXTENSIONS:
            return self._parse_image(path)
        return ParsedFile(
            name=path.name,
            content=f"지원하지 않는 파일 형식입니다: {ext}",
            mime_type="application/octet-stream",
        )

    def _parse_text(self, path: Path) -> ParsedFile:
        content = path.read_text(encoding="utf-8", errors="replace")
        mime = mimetypes.guess_type(path.name)[0] or "text/plain"
        return ParsedFile(name=path.name, content=content, mime_type=mime)

    def _parse_pdf(self, path: Path) -> ParsedFile:
        try:
            reader = PdfReader(str(path))
            pages = [page.extract_text() or "" for page in reader.pages]
            content = "\n\n".join(pages).strip()
        except Exception as e:
            content = f"PDF 파싱 오류: {e}"
        return ParsedFile(name=path.name, content=content, mime_type="application/pdf")

    def _parse_hwpx(self, path: Path) -> ParsedFile:
        try:
            from hwpx import TextExtractor

            with TextExtractor(str(path)) as extractor:
                content = extractor.extract_text(skip_empty=True)
        except ImportError:
            content = "HWPX 파싱 불가: python-hwpx 패키지가 설치되지 않았습니다."
        except Exception as e:
            content = f"HWPX 파싱 오류: {e}"
        return ParsedFile(name=path.name, content=content, mime_type="application/x-hwpx")

    def _parse_hwp(self, path: Path) -> ParsedFile:
        """구 HWP 포맷(.hwp) — pyhwp의 hwp5txt CLI를 통해 텍스트를 추출합니다."""
        try:
            result = subprocess.run(
                ["hwp5txt", str(path)],
                capture_output=True,
                text=True,
                timeout=30,
                encoding="utf-8",
                errors="replace",
            )
            if result.returncode == 0:
                content = result.stdout.strip()
            else:
                content = f"HWP 파싱 오류: {result.stderr.strip()}"
        except FileNotFoundError:
            content = "HWP 파싱 불가: pyhwp 패키지가 설치되지 않았습니다 (pip install pyhwp)."
        except subprocess.TimeoutExpired:
            content = "HWP 파싱 시간 초과 (30초)."
        except Exception as e:
            content = f"HWP 파싱 오류: {e}"
        return ParsedFile(name=path.name, content=content, mime_type="application/x-hwp")

    def _parse_docx(self, path: Path) -> ParsedFile:
        try:
            doc = Document(str(path))
            parts: list[str] = []

            # 단락
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text)

            # 표 → 마크다운 테이블
            for table in doc.tables:
                rows = [[cell.text.replace("\n", " ") for cell in row.cells] for row in table.rows]
                md_table = _rows_to_markdown(rows)
                if md_table:
                    parts.append(md_table)

            content = "\n\n".join(parts)
        except ImportError:
            content = "DOCX 파싱 불가: python-docx 패키지가 설치되지 않았습니다."
        except Exception as e:
            content = f"DOCX 파싱 오류: {e}"
        return ParsedFile(name=path.name, content=content, mime_type=DOCX_MIME)

    def _parse_xlsx(self, path: Path) -> ParsedFile:
        """Excel 파일(.xlsx/.xls) — openpyxl로 각 시트를 마크다운 테이블로 변환합니다."""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            sheets: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = [
                    [str(cell.value) if cell.value is not None else "" for cell in row]
                    for row in ws.iter_rows()
                    if any(cell.value is not None for cell in row)
                ]
                if not rows:
                    continue
                md_table = _rows_to_markdown(rows)
                sheets.append(f"### {sheet_name}\n\n{md_table}")
            wb.close()
            content = "\n\n".join(sheets) if sheets else "(빈 파일)"
        except ImportError:
            content = (
                "Excel 파싱 불가: openpyxl 패키지가 설치되지 않았습니다 (pip install openpyxl)."
            )
        except Exception as e:
            content = f"Excel 파싱 오류: {e}"
        mime = (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            if path.suffix.lower() == ".xlsx"
            else "application/vnd.ms-excel"
        )
        return ParsedFile(name=path.name, content=content, mime_type=mime)

    def _parse_pptx(self, path: Path) -> ParsedFile:
        """PowerPoint 파일(.pptx) — python-pptx로 슬라이드별 텍스트를 추출합니다."""
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            texts.append(text)
                if texts:
                    slides.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
            content = "\n\n".join(slides) if slides else "(텍스트 없음)"
        except ImportError:
            content = "PPTX 파싱 불가: python-pptx 패키지가 설치되지 않았습니다 (pip install python-pptx)."
        except Exception as e:
            content = f"PPTX 파싱 오류: {e}"
        return ParsedFile(
            name=path.name,
            content=content,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    def _parse_image(self, path: Path) -> ParsedFile:
        mime = mimetypes.guess_type(path.name)[0] or "image/jpeg"
        data = base64.b64encode(path.read_bytes()).decode("utf-8")
        return ParsedFile(name=path.name, content=data, mime_type=mime)
